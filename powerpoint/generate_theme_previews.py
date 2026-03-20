"""
PPT 테마 샘플 슬라이드 생성 + 이미지 변환 스크립트
"""
import os
import comtypes.client
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = r"D:\Claude Code\powerpoint\theme_previews"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ── 테마 정의 ─────────────────────────────────────────────
THEMES = {
    "Corporate": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x1A, 0x3A, 0x6C),   # 네이비
        "sub_accent": RGBColor(0x2E, 0x86, 0xC1), # 라이트 블루
        "text": RGBColor(0x1C, 0x1C, 0x1C),
        "sub_text": RGBColor(0x55, 0x55, 0x55),
        "bar": RGBColor(0x1A, 0x3A, 0x6C),
        "tag": "제안서 · 보고서",
    },
    "Modern": {
        "bg": RGBColor(0x12, 0x12, 0x1A),
        "accent": RGBColor(0x6C, 0x63, 0xFF),    # 퍼플
        "sub_accent": RGBColor(0x00, 0xD4, 0xAA), # 민트
        "text": RGBColor(0xF0, 0xF0, 0xF0),
        "sub_text": RGBColor(0xAA, 0xAA, 0xCC),
        "bar": RGBColor(0x6C, 0x63, 0xFF),
        "tag": "피칭 · 제안서",
    },
    "Warm": {
        "bg": RGBColor(0xFF, 0xFB, 0xF5),
        "accent": RGBColor(0xE8, 0x6B, 0x2C),    # 오렌지
        "sub_accent": RGBColor(0xF5, 0xA6, 0x23), # 옐로우
        "text": RGBColor(0x2D, 0x1B, 0x00),
        "sub_text": RGBColor(0x7A, 0x5C, 0x3A),
        "bar": RGBColor(0xE8, 0x6B, 0x2C),
        "tag": "강의 · 요약",
    },
    "Tech": {
        "bg": RGBColor(0x0D, 0x1B, 0x2A),
        "accent": RGBColor(0x00, 0xFF, 0x88),    # 그린
        "sub_accent": RGBColor(0x00, 0xB4, 0xD8), # 사이안
        "text": RGBColor(0xE0, 0xFF, 0xF0),
        "sub_text": RGBColor(0x80, 0xB3, 0xA0),
        "bar": RGBColor(0x00, 0xFF, 0x88),
        "tag": "기술문서",
    },
    "Bold": {
        "bg": RGBColor(0xF7, 0xF7, 0xF7),
        "accent": RGBColor(0xE8, 0x00, 0x54),    # 핫핑크
        "sub_accent": RGBColor(0xFF, 0xA5, 0x00), # 오렌지
        "text": RGBColor(0x0A, 0x0A, 0x0A),
        "sub_text": RGBColor(0x44, 0x44, 0x44),
        "bar": RGBColor(0xE8, 0x00, 0x54),
        "tag": "피칭",
    },
    "Clean": {
        "bg": RGBColor(0xF4, 0xF6, 0xF8),
        "accent": RGBColor(0x2C, 0x3E, 0x50),    # 다크 슬레이트
        "sub_accent": RGBColor(0x95, 0xA5, 0xA6), # 그레이
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "sub_text": RGBColor(0x7F, 0x8C, 0x8D),
        "bar": RGBColor(0x2C, 0x3E, 0x50),
        "tag": "보고서 · 요약",
    },
}

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, left, top, width, height, color, size, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def build_slide(theme_name, theme):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    W = prs.slide_width
    H = prs.slide_height

    # ── 배경 ──
    add_rect(slide, 0, 0, W, H, theme["bg"])

    # ── 좌측 사이드바 ──
    sidebar_w = Inches(3.8)
    add_rect(slide, 0, 0, sidebar_w, H, theme["accent"])

    # 사이드바 하단 서브 컬러 블록
    add_rect(slide, 0, H - Inches(1.2), sidebar_w, Inches(1.2), theme["sub_accent"])

    # ── 사이드바 텍스트 ──
    add_text(slide, theme_name, Inches(0.3), Inches(0.5), Inches(3.2), Inches(0.8),
             theme["text"] if theme["bg"] != RGBColor(0xFF,0xFF,0xFF) else RGBColor(0xFF,0xFF,0xFF),
             28, bold=True)

    # 실제 사이드바 위 텍스트는 항상 밝게
    sb_text_color = RGBColor(0xFF, 0xFF, 0xFF)
    add_text(slide, theme_name, Inches(0.3), Inches(0.5), Inches(3.2), Inches(0.8),
             sb_text_color, 28, bold=True)
    add_text(slide, theme["tag"], Inches(0.3), Inches(1.1), Inches(3.2), Inches(0.5),
             RGBColor(0xCC, 0xDD, 0xFF),
             12)

    # 메뉴 아이템들 (사이드바)
    menu_items = ["01  개요", "02  주요 내용", "03  데이터 분석", "04  결론 및 제언"]
    for i, item in enumerate(menu_items):
        y = Inches(2.0 + i * 0.7)
        if i == 1:  # 활성 메뉴
            add_rect(slide, 0, y - Inches(0.05), sidebar_w, Inches(0.55), theme["sub_accent"])
        add_text(slide, item, Inches(0.4), y, Inches(3.0), Inches(0.5),
                 RGBColor(0xFF, 0xFF, 0xFF), 11, bold=(i == 1))

    # ── 메인 영역 ──
    main_x = sidebar_w + Inches(0.3)
    content_w = W - sidebar_w - Inches(0.6)

    # 상단 타이틀 바
    add_rect(slide, main_x, Inches(0.4), content_w, Inches(0.08), theme["sub_accent"])

    add_text(slide, "슬라이드 제목이 여기에 표시됩니다", main_x, Inches(0.5),
             content_w, Inches(0.7), theme["text"], 20, bold=True)
    add_text(slide, "부제목 또는 섹션 설명 텍스트", main_x, Inches(1.15),
             content_w, Inches(0.4), theme["sub_text"], 11)

    # 구분선
    add_rect(slide, main_x, Inches(1.55), content_w, Inches(0.03), theme["sub_accent"])

    # ── 콘텐츠 카드 3개 ──
    card_w = (content_w - Inches(0.4)) / 3
    card_titles = ["핵심 지표", "분석 결과", "제언 사항"]
    card_values = ["98.5%", "▲ 24%", "3가지"]
    card_descs = ["목표 달성률", "전월 대비 성장", "개선 포인트"]

    for i in range(3):
        cx = main_x + i * (card_w + Inches(0.2))
        cy = Inches(1.75)
        ch = Inches(1.8)

        add_rect(slide, cx, cy, card_w, ch, theme["sub_accent"] if i == 1 else theme["bg"])
        # 카드 상단 컬러 바
        add_rect(slide, cx, cy, card_w, Inches(0.1), theme["accent"] if i != 1 else RGBColor(0xFF,0xFF,0xFF))

        text_color = RGBColor(0xFF,0xFF,0xFF) if i == 1 else theme["text"]
        add_text(slide, card_titles[i], cx + Inches(0.1), cy + Inches(0.15),
                 card_w - Inches(0.2), Inches(0.35), text_color, 9)
        add_text(slide, card_values[i], cx + Inches(0.1), cy + Inches(0.5),
                 card_w - Inches(0.2), Inches(0.6), text_color, 22, bold=True)
        add_text(slide, card_descs[i], cx + Inches(0.1), cy + Inches(1.1),
                 card_w - Inches(0.2), Inches(0.35), text_color, 8)

    # ── 하단 텍스트 영역 ──
    add_rect(slide, main_x, Inches(3.75), content_w, Inches(0.03), theme["sub_accent"])
    body_text = (
        "본문 텍스트는 이 영역에 배치됩니다. 데이터 기반의 인사이트와 핵심 메시지를 "
        "명확하게 전달하며, 시각적 요소와 조화를 이루는 레이아웃으로 구성됩니다."
    )
    add_text(slide, body_text, main_x, Inches(3.85), content_w, Inches(1.2),
             theme["sub_text"], 10)

    # ── 하단 바 ──
    add_rect(slide, main_x, H - Inches(0.5), content_w, Inches(0.5), theme["accent"])
    add_text(slide, "Company Name  |  2026", main_x + Inches(0.2), H - Inches(0.45),
             Inches(4), Inches(0.4), RGBColor(0xFF, 0xFF, 0xFF), 9)
    add_text(slide, "1 / 12", W - Inches(1.5), H - Inches(0.45),
             Inches(1.2), Inches(0.4), RGBColor(0xFF, 0xFF, 0xFF), 9, align=PP_ALIGN.RIGHT)

    path = os.path.join(OUTPUT_DIR, f"{theme_name}.pptx")
    prs.save(path)
    return path

# ── 생성 ──
pptx_paths = {}
for name, theme in THEMES.items():
    path = build_slide(name, theme)
    pptx_paths[name] = path
    print(f"생성: {name} → {path}")

# ── PowerPoint COM으로 PNG 변환 ──
print("\nPNG 변환 중...")
powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
powerpoint.Visible = 1

for name, pptx_path in pptx_paths.items():
    pptx_abs = os.path.abspath(pptx_path)
    png_path = os.path.join(OUTPUT_DIR, f"{name}.png")
    presentation = powerpoint.Presentations.Open(pptx_abs, ReadOnly=1, WithWindow=0)
    presentation.Slides(1).Export(os.path.abspath(png_path), "PNG", 1920, 1080)
    presentation.Close()
    print(f"PNG 저장: {png_path}")

powerpoint.Quit()
print("\n완료!")
